import os
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt

OUTPUT_DIR = "outputs"

def ensure_output_dir():
    if not os.path.exists(OUTPUT_DIR):
        os.makedirs(OUTPUT_DIR)

def to_html(report_dict):
    """
    Writes outputs/report.html showing: site name, urls crawled,
    issue counts by severity, top issues as an HTML table, recommendations list.
    """
    ensure_output_dir()
    output_path = os.path.join(OUTPUT_DIR, "report.html")

    site_name = report_dict.get("site_name", "Unknown Site")
    urls_crawled = report_dict.get("urls_crawled", 0)
    severity_counts = report_dict.get("severity_counts", {})
    top_issues = report_dict.get("top_issues", [])
    recommendations = report_dict.get("recommendations", [])

    severity_html = "".join([f"<li>{sev}: {count}</li>" for sev, count in severity_counts.items()])

    issues_rows = "".join([
        f"<tr><td>{issue.get('severity', 'N/A')}</td><td>{issue.get('issue', 'N/A')}</td><td>{issue.get('url', 'N/A')}</td></tr>"
        for issue in top_issues
    ])

    recs_html = "".join([f"<li>{rec}</li>" for rec in recommendations])

    html_content = f"""
    <!DOCTYPE html>
    <html lang="en">
    <head>
        <meta charset="UTF-8">
        <title>SEO Report - {site_name}</title>
        <style>
            body {{ font-family: Arial, sans-serif; margin: 40px; line-height: 1.6; }}
            h1, h2 {{ color: #333; }}
            table {{ border-collapse: collapse; width: 100%; margin-bottom: 20px; }}
            th, td {{ border: 1px solid #ddd; padding: 8px; text-align: left; }}
            th {{ background-color: #f2f2f2; }}
            .severity-high {{ color: red; font-weight: bold; }}
            .severity-medium {{ color: orange; font-weight: bold; }}
            .severity-low {{ color: green; font-weight: bold; }}
        </style>
    </head>
    <body>
        <h1>SEO Audit Report: {site_name}</h1>
        <p><strong>URLs Crawled:</strong> {urls_crawled}</p>

        <h2>Issue Summary</h2>
        <ul>
            {severity_html}
        </ul>

        <h2>Top Issues</h2>
        <table>
            <thead>
                <tr><th>Severity</th><th>Issue</th><th>URL</th></tr>
            </thead>
            <tbody>
                {issues_rows}
            </tbody>
        </table>

        <h2>Recommendations</h2>
        <ul>
            {recs_html}
        </ul>
    </body>
    </html>
    """

    with open(output_path, "w", encoding="utf-8") as f:
        f.write(html_content)

def to_pdf(report_dict):
    """
    Writes outputs/report.pdf using fpdf2 library.
    """
    ensure_output_dir()
    output_path = os.path.join(OUTPUT_DIR, "report.pdf")

    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", "B", 16)

    site_name = report_dict.get("site_name", "Unknown Site")
    pdf.cell(0, 10, f"SEO Audit Report: {site_name}", ln=True, align="C")
    pdf.ln(10)

    pdf.set_font("Arial", "", 12)
    urls_crawled = report_dict.get("urls_crawled", 0)
    pdf.cell(0, 10, f"URLs Crawled: {urls_crawled}", ln=True)
    pdf.ln(5)

    pdf.set_font("Arial", "B", 14)
    pdf.cell(0, 10, "Issue Summary", ln=True)
    pdf.set_font("Arial", "", 12)
    severity_counts = report_dict.get("severity_counts", {})
    for sev, count in severity_counts.items():
        pdf.cell(0, 10, f"{sev}: {count}", ln=True)
    pdf.ln(5)

    pdf.set_font("Arial", "B", 14)
    pdf.cell(0, 10, "Top Issues", ln=True)
    pdf.set_font("Arial", "", 10)

    # Simple table header
    pdf.cell(30, 10, "Severity", 1)
    pdf.cell(80, 10, "Issue", 1)
    pdf.cell(80, 10, "URL", 1)
    pdf.ln()

    top_issues = report_dict.get("top_issues", [])
    for issue in top_issues:
        pdf.cell(30, 10, issue.get("severity", "N/A"), 1)
        pdf.cell(80, 10, issue.get("issue", "N/A")[:75], 1) # Simple truncation
        pdf.cell(80, 10, issue.get("url", "N/A")[:75], 1)
        pdf.ln()
    pdf.ln(5)

    pdf.set_font("Arial", "B", 14)
    pdf.cell(0, 10, "Recommendations", ln=True)
    pdf.set_font("Arial", "", 12)
    recommendations = report_dict.get("recommendations", [])
    for rec in recommendations:
        pdf.multi_cell(0, 10, f"- {rec}")
        pdf.ln(2)

    pdf.output(output_path)

def to_pptx(report_dict):
    """
    Writes outputs/report.pptx using python-pptx library.
    """
    ensure_output_dir()
    output_path = os.path.join(OUTPUT_DIR, "report.pptx")

    prs = Presentation()

    # Slide 1: Title
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    site_name = report_dict.get("site_name", "Unknown Site")
    title.text = f"SEO Audit Report"
    subtitle.text = f"Site: {site_name}\nGenerated on 2026-06-06"

    # Slide 2: Summary
    summary_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(summary_layout)
    slide.shapes.title.text = "Executive Summary"
    tf = slide.placeholders[1].text_frame
    urls_crawled = report_dict.get("urls_crawled", 0)
    tf.text = f"Total URLs Crawled: {urls_crawled}"

    severity_counts = report_dict.get("severity_counts", {})
    for sev, count in severity_counts.items():
        p = tf.add_paragraph()
        p.text = f"{sev} Severity Issues: {count}"

    # Slide 3: Top Issues
    issues_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(issues_layout)
    slide.shapes.title.text = "Top SEO Issues"
    tf = slide.placeholders[1].text_frame
    top_issues = report_dict.get("top_issues", [])
    if not top_issues:
        tf.text = "No critical issues found."
    else:
        for issue in top_issues[:5]: # Limit to top 5 for slide readability
            p = tf.add_paragraph()
            p.text = f"[{issue.get('severity')}] {issue.get('issue')}"
            p.level = 0

    # Slide 4: Recommendations
    recs_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(recs_layout)
    slide.shapes.title.text = "Key Recommendations"
    tf = slide.placeholders[1].text_frame
    recommendations = report_dict.get("recommendations", [])
    if not recommendations:
        tf.text = "No specific recommendations at this time."
    else:
        for rec in recommendations[:5]: # Limit to top 5
            p = tf.add_paragraph()
            p.text = rec
            p.level = 0

    prs.save(output_path)
